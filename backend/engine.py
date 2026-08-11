from __future__ import annotations

import json
import re
import shutil
import time
import urllib.error
import urllib.request
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

APP_NAME = "AI-DeskClean"
APP_VERSION = "1.3.0-MVP"

OLLAMA_MODEL = "qwen2.5:1.5b"
OLLAMA_MODELS: dict[str, str] = {
    "Fast — Qwen2.5 1.5B": "qwen2.5:1.5b",
    "Lite — Qwen3 0.6B": "qwen3:0.6b",
    "Balanced — Qwen3 1.7B": "qwen2.5:1.5b",
    "Advanced — Qwen3 4B": "qwen3:4b",
}
OLLAMA_URL = "http://127.0.0.1:11434"

APP_DIR = Path.home() / ".ai_deskclean"
LOG_DIR = APP_DIR / "logs"
BACKUP_DIR = APP_DIR / "backups"
SETTINGS_PATH = APP_DIR / "settings.json"
APP_DIR.mkdir(exist_ok=True)
LOG_DIR.mkdir(exist_ok=True)
BACKUP_DIR.mkdir(exist_ok=True)

EXT_GROUPS: dict[str, set[str]] = {
    "Documents": {".pdf", ".doc", ".docx", ".txt", ".rtf", ".odt"},
    "Spreadsheets": {".xls", ".xlsx", ".csv", ".ods"},
    "Presentations": {".ppt", ".pptx", ".odp"},
    "Images": {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"},
    "Media": {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".avi", ".mkv"},
    "Archives": {".zip", ".rar", ".7z", ".tar", ".gz"},
    "Code": {".py", ".js", ".ts", ".java", ".cpp", ".c", ".html", ".css", ".json"},
}

CLASSIFICATION_METHODS: dict[str, str] = {
    "extension": "By file type (fast, offline, no AI)",
    "ai": "By academic content using a local AI model",
    "keywords": "By custom keywords (you choose the folder rules)",
    "age": "By file age",
}


@dataclass
class ScanResult:
    source: Path
    category: str
    status: str = "ready"
    selected: bool = True
    error: str = ""



def safe_name(text: str) -> str:
    text = re.sub(r"[^\w\-\u0600-\u06FF ]+", " ", text, flags=re.UNICODE)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:60] or "Other"


def read_settings() -> dict:
    try:
        return json.loads(SETTINGS_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {"language": "en", "recursive": False, "ollama_model": OLLAMA_MODEL}


def write_settings(data: dict) -> None:
    try:
        SETTINGS_PATH.write_text(
            json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    except Exception:
        pass


def extract_text(path: Path, limit: int = 6000) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".csv", ".json", ".py", ".html", ".css", ".js", ".ts"}:
            return path.read_text(encoding="utf-8", errors="ignore")[:limit]
        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            return "\n".join((p.extract_text() or "") for p in reader.pages[:6])[:limit]
        if suffix == ".docx":
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)[:limit]
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets[:4]:
                for row in ws.iter_rows(max_row=40, values_only=True):
                    parts.append(" ".join(str(v) for v in row if v is not None))
            return "\n".join(parts)[:limit]
        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts: list[str] = []
            for slide in prs.slides[:12]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)[:limit]
    except Exception:
        return ""
    return ""


def extension_category(path: Path) -> str:
    for name, exts in EXT_GROUPS.items():
        if path.suffix.lower() in exts:
            return name
    return "Other"


def age_category(path: Path) -> str:
    age_days = (time.time() - path.stat().st_mtime) / 86400
    if age_days <= 30:
        return "Last 30 Days"
    if age_days <= 180:
        return "1-6 Months"
    if age_days <= 365:
        return "6-12 Months"
    if age_days <= 730:
        return "1-2 Years"
    return "Older than 2 Years"


def parse_rules(text: str) -> list[tuple[str, list[str]]]:
    """Parse the raw 'Folder = kw1, kw2' textarea format."""
    rules: list[tuple[str, list[str]]] = []
    for line in text.splitlines():
        if "=" not in line:
            continue
        folder, words = line.split("=", 1)
        keywords = [w.strip().lower() for w in re.split(r"[,\u060c]", words) if w.strip()]
        if folder.strip() and keywords:
            rules.append((safe_name(folder.strip()), keywords))
    return rules


def rules_from_structured(rules: list[dict[str, Any]] | None) -> list[tuple[str, list[str]]]:
    """
    Same normalization as parse_rules(), but starting from the
    already-structured ``[{"folder": ..., "keywords": [...]}, ...]``
    list the HTML frontend sends (it parses its own textarea client
    side before calling the bridge).
    """
    result: list[tuple[str, list[str]]] = []
    for rule in rules or []:
        folder = str(rule.get("folder", "")).strip()
        keywords = [str(k).strip().lower() for k in (rule.get("keywords") or []) if str(k).strip()]
        if folder and keywords:
            result.append((safe_name(folder), keywords))
    return result


def rule_category(path: Path, rules: list[tuple[str, list[str]]], include_content: bool) -> str:
    haystack = path.stem.lower()
    if include_content:
        haystack += " " + extract_text(path, 2500).lower()
    for category, words in rules:
        if any(word in haystack for word in words):
            return category
    return "Other"


def ollama_ready(model: str = OLLAMA_MODEL) -> bool:
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=3) as response:
            data = json.loads(response.read().decode("utf-8"))
        names = {m.get("name", "") for m in data.get("models", [])}
        return any(name == model or name.startswith(model + ":") for name in names)
    except Exception:
        return False


def normalize_ai_category(value: str, fallback: str = "General Academic") -> str:
    """Normalize an AI-proposed folder name into a short, safe, useful category."""
    value = safe_name(value)
    aliases = {
        "lecture": "Lectures", "lectures": "Lectures", "\u0645\u062d\u0627\u0636\u0631\u0629": "\u0645\u062d\u0627\u0636\u0631\u0627\u062a", "\u0645\u062d\u0627\u0636\u0631\u0627\u062a": "\u0645\u062d\u0627\u0636\u0631\u0627\u062a",
        "assignment": "Assignments", "assignments": "Assignments", "homework": "Assignments", "\u0648\u0627\u062c\u0628": "\u0648\u0627\u062c\u0628\u0627\u062a", "\u0648\u0627\u062c\u0628\u0627\u062a": "\u0648\u0627\u062c\u0628\u0627\u062a",
        "project": "Projects", "projects": "Projects", "\u0645\u0634\u0631\u0648\u0639": "\u0645\u0634\u0627\u0631\u064a\u0639", "\u0645\u0634\u0627\u0631\u064a\u0639": "\u0645\u0634\u0627\u0631\u064a\u0639",
        "research": "Research", "paper": "Research", "papers": "Research", "\u0628\u062d\u062b": "\u0623\u0628\u062d\u0627\u062b", "\u0623\u0628\u062d\u0627\u062b": "\u0623\u0628\u062d\u0627\u062b",
        "artificial intelligence": "Artificial Intelligence", "ai": "Artificial Intelligence", "\u0630\u0643\u0627\u0621 \u0627\u0635\u0637\u0646\u0627\u0639\u064a": "\u0627\u0644\u0630\u0643\u0627\u0621 \u0627\u0644\u0627\u0635\u0637\u0646\u0627\u0639\u064a",
        "cybersecurity": "Cybersecurity", "cyber security": "Cybersecurity", "\u0623\u0645\u0646 \u0633\u064a\u0628\u0631\u0627\u0646\u064a": "\u0627\u0644\u0623\u0645\u0646 \u0627\u0644\u0633\u064a\u0628\u0631\u0627\u0646\u064a",
        "programming": "Programming", "code": "Programming", "\u0628\u0631\u0645\u062c\u0629": "\u0627\u0644\u0628\u0631\u0645\u062c\u0629",
        "database": "Databases", "databases": "Databases", "\u0642\u0648\u0627\u0639\u062f \u0628\u064a\u0627\u0646\u0627\u062a": "\u0642\u0648\u0627\u0639\u062f \u0627\u0644\u0628\u064a\u0627\u0646\u0627\u062a",
        "software engineering": "Software Engineering", "\u0647\u0646\u062f\u0633\u0629 \u0628\u0631\u0645\u062c\u064a\u0627\u062a": "\u0647\u0646\u062f\u0633\u0629 \u0627\u0644\u0628\u0631\u0645\u062c\u064a\u0627\u062a",
        "mathematics": "Mathematics", "math": "Mathematics", "\u0631\u064a\u0627\u0636\u064a\u0627\u062a": "\u0627\u0644\u0631\u064a\u0627\u0636\u064a\u0627\u062a",
        "administration": "Administration", "admin": "Administration", "\u0625\u062f\u0627\u0631\u0629": "\u0627\u0644\u0625\u062f\u0627\u0631\u0629",
        "cv and career": "CV and Career", "career": "CV and Career", "resume": "CV and Career", "\u0633\u064a\u0631\u0629 \u0630\u0627\u062a\u064a\u0629": "\u0627\u0644\u0633\u064a\u0631\u0629 \u0627\u0644\u0630\u0627\u062a\u064a\u0629 \u0648\u0627\u0644\u062a\u0648\u0638\u064a\u0641",
        "finance": "Finance", "invoice": "Finance", "\u0641\u0648\u0627\u062a\u064a\u0631": "\u0627\u0644\u0645\u0627\u0644\u064a\u0629 \u0648\u0627\u0644\u0641\u0648\u0627\u062a\u064a\u0631",
        "design and media": "Design and Media", "design": "Design and Media", "media": "Design and Media", "\u062a\u0635\u0645\u064a\u0645": "\u0627\u0644\u062a\u0635\u0645\u064a\u0645 \u0648\u0627\u0644\u0648\u0633\u0627\u0626\u0637",
        "general academic": "General Academic", "academic": "General Academic", "\u0623\u0643\u0627\u062f\u064a\u0645\u064a \u0639\u0627\u0645": "\u0623\u0643\u0627\u062f\u064a\u0645\u064a \u0639\u0627\u0645",
        "other": "Other", "\u0623\u062e\u0631\u0649": "\u0623\u062e\u0631\u0649",
    }
    return aliases.get(value.lower(), value or fallback)


def academic_category_with_confidence(path: Path, model: str = OLLAMA_MODEL) -> tuple[str, float]:
    """Classify a file with the local AI model and return its confidence."""
    text = extract_text(path, 9000)
    ext_hint = extension_category(path)
    if not text.strip():
        return normalize_ai_category(ext_hint, "Other"), 0.0

    prompt = f"""You are the classification engine of a local desktop file organizer.
Analyze BOTH the file name and the extracted document content, then choose the most useful destination folder.

Rules:
1. Return strict JSON only: {{"category":"...","confidence":0.0,"reason":"..."}}
2. category must be a short folder name of 1 to 4 words, not a sentence.
3. Prefer a specific academic topic when the content clearly supports it.
4. If the file is a lecture, assignment, project, or research paper, prefer that document-purpose category unless a strong subject category is more useful.
5. Do not invent a category from a single weak word. Use the dominant repeated topic and headings.
6. Support Arabic and English. Use the dominant language of the document for the folder name.
7. Suitable categories include: Lectures, Assignments, Projects, Research, Artificial Intelligence, Cybersecurity, Programming, Databases, Software Engineering, Mathematics, Administration, CV and Career, Finance, Design and Media, General Academic, Other.
8. confidence must be between 0 and 1.

File name: {path.name}
File extension group: {ext_hint}
Extracted content:
{text[:9000]}
"""
    payload = json.dumps({
        "model": model,
        "messages": [
            {"role": "system", "content": "You classify local files. Output valid JSON only."},
            {"role": "user", "content": prompt},
        ],
        "format": "json",
        "stream": False,
        "options": {
            "temperature": 0,
            "num_ctx": 4096,
            "num_predict": 100,
            "top_p": 0.8,
        },
    }).encode("utf-8")
    request = urllib.request.Request(
        f"{OLLAMA_URL}/api/chat",
        data=payload,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=120) as response:
        result = json.loads(response.read().decode("utf-8"))
    content = result.get("message", {}).get("content", "{}")
    try:
        parsed = json.loads(content)
        category = normalize_ai_category(str(parsed.get("category", "General Academic")))
        confidence = float(parsed.get("confidence", 0.0) or 0.0)
    except (ValueError, TypeError, json.JSONDecodeError):
        return normalize_ai_category(ext_hint, "Other"), 0.0

    if confidence < 0.35 or category.lower() in {"other", "\u0623\u062e\u0631\u0649"}:
        return normalize_ai_category(ext_hint, "Other"), confidence
    return category, confidence


def academic_category(path: Path, model: str = OLLAMA_MODEL) -> str:
    """Classify a file with the local AI model."""
    category, _confidence = academic_category_with_confidence(path, model)
    return category


def list_operations() -> list[dict]:
    """Load every saved operation manifest from the backups folder, newest first."""
    operations = []
    for manifest_path in BACKUP_DIR.glob("operation-*.json"):
        try:
            data = json.loads(manifest_path.read_text(encoding="utf-8"))
            data["_path"] = str(manifest_path)
            operations.append(data)
        except Exception:
            continue
    operations.sort(key=lambda item: item.get("created_at", ""), reverse=True)
    return operations


def compute_dashboard_stats(operations: list[dict]) -> dict:
    """Turn the raw manifest list into summary numbers (category/method counts, etc.)."""
    stats = {
        "files_organized": 0,
        "operations": 0,
        "categories": 0,
        "undone": 0,
        "last_run": None,
    }
    categories = set()
    for op in operations:
        moves = op.get("moves", [])
        stats["operations"] += 1
        if op.get("undone"):
            stats["undone"] += 1
        else:
            stats["files_organized"] += len(moves)
            for move in moves:
                categories.add(Path(move["destination"]).parent.name)
    stats["categories"] = len(categories)
    if operations:
        stats["last_run"] = operations[0].get("created_at")
    return stats


def report_breakdown(operations: list[dict]) -> dict[str, Any]:
    """
    Aggregate cumulative report data across every organize operation,
    ported from the original Tkinter build's show_reports()/_report_panel()
    logic (category counts + method counts over every non-undone run).
    """
    active_ops = [op for op in operations if not op.get("undone")]

    category_counts: dict[str, int] = {}
    method_counts: dict[str, int] = {}

    for op in active_ops:
        method = op.get("method") or "extension"
        moves = op.get("moves", [])
        method_counts[method] = method_counts.get(method, 0) + len(moves)
        for move in moves:
            cat = Path(move["destination"]).parent.name
            category_counts[cat] = category_counts.get(cat, 0) + 1

    by_category = sorted(category_counts.items(), key=lambda kv: kv[1], reverse=True)
    by_method = sorted(method_counts.items(), key=lambda kv: kv[1], reverse=True)

    total_files = sum(category_counts.values())
    total_operations = len(active_ops)

    return {
        "has_data": bool(active_ops),
        "total_files": total_files,
        "total_operations": total_operations,
        "by_category": [{"name": name, "count": count} for name, count in by_category],
        "by_method": [{"method": method, "count": count} for method, count in by_method],
    }


def unique_destination(folder: Path, name: str) -> Path:
    candidate = folder / name
    if not candidate.exists():
        return candidate
    stem, suffix = Path(name).stem, Path(name).suffix
    for i in range(1, 10000):
        candidate = folder / f"{stem} ({i}){suffix}"
        if not candidate.exists():
            return candidate
    raise RuntimeError("Could not create a unique destination name")



def collect_files(folder: Path, recursive: bool) -> list[Path]:
    """Collect files that can be organized."""
    if not folder:
        return []
    iterator = folder.rglob("*") if recursive else folder.iterdir()
    files: list[Path] = []
    for p in iterator:
        try:
            if not p.is_file():
                continue
            if p.name.startswith("~$") or p.name.startswith("."):
                continue
            if APP_DIR in p.parents:
                continue
            if p.suffix.lower() in {".lnk", ".ini", ".sys", ".dll", ".exe"}:
                continue
            files.append(p)
        except OSError:
            continue
    return files


def move_files(folder: Path, chosen: list[dict[str, Any]]) -> tuple[dict, float]:
    """Move the selected files and record the operation."""
    start = time.perf_counter()
    moves: list[dict[str, str]] = []
    errors: list[dict[str, str]] = []

    for item in chosen:
        source = Path(item["path"])
        try:
            dest_dir = folder / safe_name(item["category"])
            dest_dir.mkdir(exist_ok=True)
            dest = unique_destination(dest_dir, source.name)
            shutil.move(str(source), str(dest))
            moves.append({"source": str(source), "destination": str(dest)})
        except Exception as exc:
            errors.append({"file": str(source), "error": str(exc)})

    time_taken = time.perf_counter() - start

    manifest = {
        "created_at": datetime.now().isoformat(),
        "root": str(folder),
        "method": None,
        "model": None,
        "undone": False,
        "moves": moves,
        "errors": errors,
    }
    return manifest, time_taken


def save_manifest(manifest: dict) -> Path:
    """Save the operation manifest and log."""
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    manifest_path = BACKUP_DIR / f"operation-{timestamp}.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (BACKUP_DIR / "latest.txt").write_text(str(manifest_path), encoding="utf-8")
    log_path = LOG_DIR / f"log-{timestamp}.json"
    log_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return manifest_path


def undo_manifest_path(manifest_path: Path) -> dict[str, Any]:
    """Restore files from an operation manifest."""
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except Exception:
        return {
            "success": False,
            "restored_count": 0,
            "error_count": 0,
            "errors": [],
            "message": "There is no previous operation to undo.",
        }

    if manifest.get("undone"):
        return {
            "success": False,
            "restored_count": 0,
            "error_count": 0,
            "errors": [],
            "message": "There is no previous operation to undo.",
        }

    restored = 0
    undo_errors: list[dict[str, str]] = []

    for move in reversed(manifest.get("moves", [])):
        src = Path(move["destination"])
        dest = Path(move["source"])
        if not src.exists():
            undo_errors.append({"file": str(src), "error": "Moved file no longer exists"})
            continue
        try:
            dest.parent.mkdir(parents=True, exist_ok=True)
            if dest.exists():
                dest = unique_destination(dest.parent, dest.name)
            shutil.move(str(src), str(dest))
            restored += 1
            try:
                if src.parent.exists() and not any(src.parent.iterdir()):
                    src.parent.rmdir()
            except OSError:
                pass
        except Exception as exc:
            undo_errors.append({"file": str(src), "error": str(exc)})

    manifest["undone"] = restored > 0 and not undo_errors
    manifest["undo_at"] = datetime.now().isoformat()
    manifest["undo_errors"] = undo_errors
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    if manifest["undone"]:
        latest = BACKUP_DIR / "latest.txt"
        try:
            if latest.exists() and latest.read_text(encoding="utf-8").strip() == str(manifest_path):
                latest.unlink(missing_ok=True)
        except Exception:
            pass

    return {
        "success": bool(manifest["undone"]),
        "restored_count": restored,
        "error_count": len(undo_errors),
        "errors": [{"path": e["file"], "error": e["error"]} for e in undo_errors],
        "message": "Undo completed." if manifest["undone"] else "Undo completed with errors.",
    }
